import os
from dotenv import load_dotenv
import streamlit as st
from groq import Groq
from pptx import Presentation
from io import BytesIO

# Load environment variables from .env file
load_dotenv()

# Set up Groq client
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
client = Groq(api_key=GROQ_API_KEY)



def gen_data(presentation_topic,selected_slide_types):

    if presentation_topic and selected_slide_types:
        # Generate prompt for Groq API with the requested structure
        prompt = f"Create a presentation on the topic '{presentation_topic}' with the following slides:\n"
        
        # Add each selected slide type to the prompt in the desired format
        for slide_type in selected_slide_types:
            prompt += f"\n- Slide Title: {slide_type}\n  Content: Please generate content for the {slide_type} slide.\n"
        
        # Call the Groq API to generate presentation content
        try:
            chat_completion = client.chat.completions.create(
                messages=[
                    {"role": "user", "content": prompt}
                ],
                model="llama3-8b-8192",
            )
            
            # Get the generated text
            generated_text = chat_completion.choices[0].message.content

            return generated_text
        except Exception as e:
            return "soemthing wnet wrong"

def gen_ppt(slide_content):
    # Create a PowerPoint presentation
    prs = Presentation()

    for slide_title, slide_text in slide_content.items():
        # Add a slide with a title and content layout
        slide_layout = prs.slide_layouts[1]  # Use layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide title
        title = slide.shapes.title
        title.text = slide_title
        
        # Add slide content
        content = slide.shapes.placeholders[1]
        content.text = slide_text

    # Save the presentation to a BytesIO object for download
    ppt_file = BytesIO()
    prs.save(ppt_file)
    ppt_file.seek(0)
    return ppt_file

# Define slide types
slide_types = [
    "Introduction", "Summary", "Advantages", "Disadvantages", "Architecture", 
    "Process", "Conclusion", "Index", "Agenda", "Challenges", "Key Findings", 
    "Market Analysis", "Results", "Methods", "Objective", "Timeline", 
    "Team", "Budget", "References", "Acknowledgments", "Closing"
]

# Streamlit UI
st.set_page_config(
    page_title="PPT Creation", 
    layout='wide',
    page_icon="📊"  # Custom page icon
)

col1, col2 = st.columns([1, 1])


with col1:

    st.title("PPT Creation with AI")
    # Input for presentation topic
    presentation_topic = st.text_input("Enter your presentation topic:")

    # Multi-select for slide types
    selected_slide_types = st.multiselect(
        "Select slide types:",
        slide_types,
        default=["Introduction", "Summary"]
    )

    # Button to generate presentation text
    if st.button("Generate Presentation Text"):

        st.subheader("Generated Presentation Content:")
        generated_text=gen_data(presentation_topic,selected_slide_types)
        # st.write(generated_text)
        st.markdown(
            f'<div style="height: 300px; overflow-y: scroll; padding: 10px; border: 1px solid #ccc; background-color: #f9f9f9; border-radius: 5px;">{generated_text}</div>',
            unsafe_allow_html=True
        )

    else:
        st.error("Please enter a topic and select slide types.")




with col2:
    st.subheader("Paste Content Into Corresponding Textboxes:")

    slide_text_boxes = {}  # Store content locally for this form session
    with st.form(key="slide_form"):
        # Display text boxes for the selected slide types

         # Limit the height of the form and make it scrollable
       
        st.subheader("Customize Slide Content:")
                
        for slide_type in selected_slide_types:
                        # Create or update text box content from local variable
            slide_text_boxes[slide_type] = st.text_area(
                            f"Slide - {slide_type}",
                            value="",
                            placeholder=f"Enter content for {slide_type} slide here...",
                            height=68,
                            key=f"slide_{slide_type}"
                        )

                # Submit button to show content or generate PPT
        submit_button = st.form_submit_button(label="Generate  PPT")

    if submit_button:
                    
            
            # st.write(slide_text_boxes)
        ppt_file=gen_ppt(slide_text_boxes)
        st.download_button(
                label="Download PPT",
                data=ppt_file,
                file_name="presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
            st.error('hit button to downlaod ppt')
